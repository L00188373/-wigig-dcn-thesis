#!/usr/bin/env python3
"""
Task 4 PowerPoint Presentation Generator
Creates presentation answering Saim's questions about packet drops
Author: Anthony Malone (L00188373)
Date: November 24, 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Background color
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(46, 134, 171)  # Blue

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Task 4: Packet Drop Analysis"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(48)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Investigating Multi-Path TCP Reordering in Hybrid DCN"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
subtitle_para.alignment = PP_ALIGN.CENTER

# Author info
author_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
author_frame = author_box.text_frame
author_frame.text = "Anthony Malone (L00188373)\nSupervisor: Dr. Saim Ghafoor\nNovember 24, 2025"
author_para = author_frame.paragraphs[0]
author_para.font.size = Pt(18)
author_para.font.color.rgb = RGBColor(255, 255, 255)
author_para.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 2: The Question
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
title2 = slide2.shapes.title
title2.text = "Saim's Question"

content2 = slide2.placeholders[1]
tf2 = content2.text_frame
tf2.text = "Why are 760 packets dropped with wireless enabled compared to 12 without wireless when there's no congestion?"

p2 = tf2.add_paragraph()
p2.text = "\nKey Points to Investigate:"
p2.level = 0
p2.font.bold = True

points = [
    "What type of packets are being dropped?",
    "Where are packets being dropped?",
    "Why does adding capacity make things worse?",
    "Is this an ns-3 limitation or real-world phenomenon?"
]

for point in points:
    p = tf2.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

# ============================================================================
# SLIDE 3: Simulation Results - The Paradox
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "Simulation Results: The Paradox"

content3 = slide3.placeholders[1]
tf3 = content3.text_frame
tf3.text = "Configuration: 4 ToR + 2 Spine, 80 servers, Uniform traffic, No congestion"

# Add table
rows = 5
cols = 4
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(3)

table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

# Header row
headers = ['Metric', 'Wired-Only', 'Hybrid (+WiGig)', 'Change']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(16)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(46, 134, 171)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
data = [
    ['Packet Drops', '12', '760', '+6,233% ❌'],
    ['Throughput', '127.2 Mbps', '115.8 Mbps', '-9.0% ❌'],
    ['Latency', '2.64 ms', '2.91 ms', '+10.2% ❌'],
    ['Packet Loss', '0.0004%', '0.023%', '+58× ❌']
]

for i, row_data in enumerate(data, start=1):
    for j, cell_data in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_data
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if '❌' in cell_data:
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            cell.text_frame.paragraphs[0].font.bold = True

# ============================================================================
# SLIDE 4: Graph - Packet Drops
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame4 = title_box4.text_frame
title_frame4.text = "Packet Drops: 63× Increase with Wireless"
title_para4 = title_frame4.paragraphs[0]
title_para4.font.size = Pt(32)
title_para4.font.bold = True
title_para4.alignment = PP_ALIGN.CENTER

# Add graph
img_path4 = 'Graphs/task4-packet-drops.png'
left4 = Inches(1)
top4 = Inches(1.2)
height4 = Inches(5.5)
slide4.shapes.add_picture(img_path4, left4, top4, height=height4)

# ============================================================================
# SLIDE 5: Graph - Performance Metrics
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame5 = title_box5.text_frame
title_frame5.text = "All Performance Metrics Degrade"
title_para5 = title_frame5.paragraphs[0]
title_para5.font.size = Pt(32)
title_para5.font.bold = True
title_para5.alignment = PP_ALIGN.CENTER

img_path5 = 'Graphs/task4-performance-metrics.png'
left5 = Inches(0.5)
top5 = Inches(1.2)
height5 = Inches(5.5)
slide5.shapes.add_picture(img_path5, left5, top5, height=height5)

# ============================================================================
# SLIDE 6: Root Cause Explanation
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
title6.text = "Root Cause: Multi-Path TCP Reordering"

content6 = slide6.placeholders[1]
tf6 = content6.text_frame
tf6.text = "NOT actual packet loss - TCP retransmissions!"

p6_1 = tf6.add_paragraph()
p6_1.text = "\nWhat Happens:"
p6_1.level = 0
p6_1.font.bold = True
p6_1.font.size = Pt(22)

steps = [
    "ECMP routing creates multiple paths (wired 3 hops, wireless 4 hops)",
    "Different path lengths → Different latencies",
    "Packets from SAME flow take DIFFERENT paths",
    "Fast packets overtake slow packets → OUT OF ORDER arrival",
    "TCP receiver detects gaps → Sends duplicate ACKs",
    "Sender interprets as loss → Fast retransmit triggered",
    "Congestion window collapses → Throughput drops 9%"
]

for step in steps:
    p = tf6.add_paragraph()
    p.text = step
    p.level = 1
    p.font.size = Pt(16)

# ============================================================================
# SLIDE 7: Visual Explanation
# ============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_box7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame7 = title_box7.text_frame
title_frame7.text = "Multi-Path Problem Visualized"
title_para7 = title_frame7.paragraphs[0]
title_para7.font.size = Pt(32)
title_para7.font.bold = True
title_para7.alignment = PP_ALIGN.CENTER

img_path7 = 'Graphs/task4-multipath-problem.png'
left7 = Inches(0.5)
top7 = Inches(1.2)
height7 = Inches(5.8)
slide7.shapes.add_picture(img_path7, left7, top7, height=height7)

# ============================================================================
# SLIDE 8: Literature Validation
# ============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title8 = slide8.shapes.title
title8.text = "Literature Validation: NOT an ns-3 Bug!"

content8 = slide8.placeholders[1]
tf8 = content8.text_frame
tf8.text = "Sur et al. (2017) - MobiCom '17"
tf8.paragraphs[0].font.bold = True
tf8.paragraphs[0].font.size = Pt(24)
tf8.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)

p8_1 = tf8.add_paragraph()
p8_1.text = '\n"WiFi-Assisted 60 GHz Wireless Networks"'
p8_1.font.italic = True
p8_1.font.size = Pt(20)

p8_2 = tf8.add_paragraph()
p8_2.text = "\nKey Findings:"
p8_2.font.bold = True
p8_2.font.size = Pt(22)

findings = [
    "5× throughput reduction with dual 60GHz+WiFi interfaces",
    "Root cause: 50% out-of-order packets",
    "TCP congestion window collapse observed",
    "Tested on REAL commodity IEEE 802.11ad hardware",
    "Published in top-tier ACM conference"
]

for finding in findings:
    p = tf8.add_paragraph()
    p.text = finding
    p.level = 1
    p.font.size = Pt(18)

p8_final = tf8.add_paragraph()
p8_final.text = "\n✓ Our simulation reproduces their real-world findings!"
p8_final.font.bold = True
p8_final.font.size = Pt(20)
p8_final.font.color.rgb = RGBColor(0, 128, 0)

# ============================================================================
# SLIDE 9: Answer to Question
# ============================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title9 = slide9.shapes.title
title9.text = "Answer to Saim's Question"

content9 = slide9.placeholders[1]
tf9 = content9.text_frame
tf9.text = 'Q: "Why 760 drops vs 12 with no congestion?"'
tf9.paragraphs[0].font.bold = True
tf9.paragraphs[0].font.size = Pt(22)

p9_1 = tf9.add_paragraph()
p9_1.text = "\nA: TCP packet reordering, NOT actual loss"
p9_1.font.bold = True
p9_1.font.size = Pt(24)
p9_1.font.color.rgb = RGBColor(0, 128, 0)

p9_2 = tf9.add_paragraph()
p9_2.text = "\nWhat type of packets?"
p9_2.level = 0
p9_2.font.bold = True

p9_3 = tf9.add_paragraph()
p9_3.text = "TCP retransmissions (FlowMonitor counts as drops)"
p9_3.level = 1

p9_4 = tf9.add_paragraph()
p9_4.text = "\nWhere dropped?"
p9_4.level = 0
p9_4.font.bold = True

p9_5 = tf9.add_paragraph()
p9_5.text = "NOT dropped - packets arrive out-of-order at receiver"
p9_5.level = 1

p9_6 = tf9.add_paragraph()
p9_6.text = "\nReal-world or simulation artifact?"
p9_6.level = 0
p9_6.font.bold = True

p9_7 = tf9.add_paragraph()
p9_7.text = "REAL phenomenon (validated by Sur et al. 2017)"
p9_7.level = 1
p9_7.font.color.rgb = RGBColor(255, 0, 0)

# ============================================================================
# SLIDE 10: Research Implications
# ============================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
title10 = slide10.shapes.title
title10.text = "Implications for Thesis"

content10 = slide10.placeholders[1]
tf10 = content10.text_frame
tf10.text = "✓ Validates Hypothesis H2"
tf10.paragraphs[0].font.bold = True
tf10.paragraphs[0].font.size = Pt(24)
tf10.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)

p10_1 = tf10.add_paragraph()
p10_1.text = '"Hybrid architectures only help when there\'s actual congestion to relieve"'
p10_1.level = 1
p10_1.font.italic = True

p10_2 = tf10.add_paragraph()
p10_2.text = "\n✓ Answers Research Question RQ2"
p10_2.font.bold = True
p10_2.font.size = Pt(24)
p10_2.font.color.rgb = RGBColor(0, 128, 0)

p10_3 = tf10.add_paragraph()
p10_3.text = "When does WiGig help? ONLY with proper traffic engineering!"
p10_3.level = 1

p10_4 = tf10.add_paragraph()
p10_4.text = "\nKey Insights:"
p10_4.font.bold = True
p10_4.font.size = Pt(22)

insights = [
    "Adding capacity ≠ Better performance",
    "Multi-path requires per-flow routing (not per-packet ECMP)",
    "Wireless needs traffic engineering, not blanket deployment",
    "Critical analysis: When NOT to use wireless"
]

for insight in insights:
    p = tf10.add_paragraph()
    p.text = insight
    p.level = 1
    p.font.size = Pt(18)

# ============================================================================
# SLIDE 11: Next Steps
# ============================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[1])
title11 = slide11.shapes.title
title11.text = "Next Steps for Task 4"

content11 = slide11.placeholders[1]
tf11 = content11.text_frame
tf11.text = "Current Status: Understood when wireless DOESN'T help"

p11_1 = tf11.add_paragraph()
p11_1.text = "\nNext: Demonstrate when it DOES help"
p11_1.font.bold = True
p11_1.font.size = Pt(22)

steps_next = [
    "Fix topology: WiGig only between ToR 2-3 (not all racks)",
    "Create congestion: Hotspot traffic Rack 3 → Rack 1",
    "Implement mice/elephant flow distribution (80%/20%)",
    "Add packet drop instrumentation by flow type",
    "Measure congestion relief with WiGig offload",
    "Compare: <60% load (no benefit) vs >80% load (benefit expected)"
]

for step in steps_next:
    p = tf11.add_paragraph()
    p.text = step
    p.level = 1
    p.font.size = Pt(18)

p11_final = tf11.add_paragraph()
p11_final.text = "\nTimeline: 2 weeks for complete Task 4 analysis"
p11_final.font.bold = True
p11_final.font.size = Pt(20)

# ============================================================================
# SLIDE 12: Summary
# ============================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
background12 = slide12.background
fill12 = background12.fill
fill12.solid()
fill12.fore_color.rgb = RGBColor(46, 134, 171)

summary_box = slide12.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5.5))
summary_frame = summary_box.text_frame
summary_frame.text = "Summary"
summary_frame.paragraphs[0].font.size = Pt(44)
summary_frame.paragraphs[0].font.bold = True
summary_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
summary_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

summary_points = [
    "",
    "✓ 760 drops caused by TCP reordering, not loss",
    "",
    "✓ Validated by Sur et al. (2017) real-world study",
    "",
    "✓ Proves H2: Hybrid only helps WITH congestion",
    "",
    "✓ Critical finding: When NOT to use wireless",
    "",
    "✓ Next: Show when wireless DOES provide benefit"
]

for point in summary_points:
    p = summary_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 13: Thank You
# ============================================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
background13 = slide13.background
fill13 = background13.fill
fill13.solid()
fill13.fore_color.rgb = RGBColor(46, 134, 171)

thanks_box = slide13.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "Questions?"
thanks_frame.paragraphs[0].font.size = Pt(64)
thanks_frame.paragraphs[0].font.bold = True
thanks_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

contact_box = slide13.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
contact_frame = contact_box.text_frame
contact_frame.text = "Anthony Malone (L00188373)\nEmail: L00188373@atu.ie"
contact_frame.paragraphs[0].font.size = Pt(20)
contact_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('Task4_PacketDrop_Analysis_Presentation.pptx')
print("✓ PowerPoint presentation created successfully!")
print("  File: Task4_PacketDrop_Analysis_Presentation.pptx")
print("  Slides: 13")
print("  Ready for Friday meeting! 🎯")
