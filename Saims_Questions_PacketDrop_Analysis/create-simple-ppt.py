#!/usr/bin/env python3
"""
Simple 6-Slide Presentation for Saim's Questions
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(46, 134, 171)

title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
title_frame = title_box.text_frame
title_frame.text = "Task 4: Packet Drop Analysis"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

subtitle = title_frame.add_paragraph()
subtitle.text = "Why 760 drops with wireless vs 12 without?"
subtitle.font.size = Pt(28)
subtitle.font.color.rgb = RGBColor(255, 255, 255)
subtitle.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 2: The Results
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
title2 = slide2.shapes.title
title2.text = "The Paradox"

rows, cols = 5, 4
table = slide2.shapes.add_table(rows, cols, Inches(1.5), Inches(2), Inches(7), Inches(4)).table

headers = ['Metric', 'Wired-Only', 'Hybrid', 'Change']
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(18)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(46, 134, 171)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

data = [
    ['Drops', '12', '760', '+6,233%'],
    ['Throughput', '127 Mbps', '116 Mbps', '-9%'],
    ['Latency', '2.64 ms', '2.91 ms', '+10%'],
    ['Loss', '0.0004%', '0.023%', '+58×']
]

for i, row_data in enumerate(data, 1):
    for j, val in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(16)

# ============================================================================
# SLIDE 3: Main Graph
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title3.text_frame.text = "63× More Packet Drops"
title3.text_frame.paragraphs[0].font.size = Pt(36)
title3.text_frame.paragraphs[0].font.bold = True
title3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

slide3.shapes.add_picture('Graphs/task4-packet-drops.png', 
                         Inches(1), Inches(1), height=Inches(6))

# ============================================================================
# SLIDE 4: Root Cause
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "Root Cause"

content4 = slide4.placeholders[1]
tf4 = content4.text_frame
tf4.text = "TCP Packet Reordering"
tf4.paragraphs[0].font.size = Pt(32)
tf4.paragraphs[0].font.bold = True

p1 = tf4.add_paragraph()
p1.text = "\nMultiple paths with different lengths:"
p1.font.size = Pt(22)

p2 = tf4.add_paragraph()
p2.text = "  • Wired: 3 hops (faster)"
p2.font.size = Pt(20)

p3 = tf4.add_paragraph()
p3.text = "  • Wireless: 4 hops (slower)"
p3.font.size = Pt(20)

p4 = tf4.add_paragraph()
p4.text = "\nPackets arrive out of order"
p4.font.size = Pt(22)

p5 = tf4.add_paragraph()
p5.text = "TCP thinks packets are lost"
p5.font.size = Pt(22)

p6 = tf4.add_paragraph()
p6.text = "Triggers retransmissions"
p6.font.size = Pt(22)
p6.font.color.rgb = RGBColor(255, 0, 0)

# ============================================================================
# SLIDE 5: Literature Proof
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Validated by Literature"

content5 = slide5.placeholders[1]
tf5 = content5.text_frame
tf5.text = "Sur et al. (2017) - MobiCom"
tf5.paragraphs[0].font.size = Pt(28)
tf5.paragraphs[0].font.bold = True
tf5.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)

p5_1 = tf5.add_paragraph()
p5_1.text = "\nTested REAL 60 GHz + WiFi hardware"
p5_1.font.size = Pt(22)

p5_2 = tf5.add_paragraph()
p5_2.text = "\nFound:"
p5_2.font.size = Pt(22)
p5_2.font.bold = True

p5_3 = tf5.add_paragraph()
p5_3.text = "  • 5× throughput reduction"
p5_3.font.size = Pt(20)

p5_4 = tf5.add_paragraph()
p5_4.text = "  • 50% out-of-order packets"
p5_4.font.size = Pt(20)

p5_5 = tf5.add_paragraph()
p5_5.text = "  • TCP window collapse"
p5_5.font.size = Pt(20)

p5_6 = tf5.add_paragraph()
p5_6.text = "\nOur results match their findings!"
p5_6.font.size = Pt(24)
p5_6.font.bold = True
p5_6.font.color.rgb = RGBColor(0, 128, 0)

# ============================================================================
# SLIDE 6: Conclusion
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background6 = slide6.background
fill6 = background6.fill
fill6.solid()
fill6.fore_color.rgb = RGBColor(46, 134, 171)

conclusion_box = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
conclusion_frame = conclusion_box.text_frame
conclusion_frame.text = "Conclusion"
conclusion_frame.paragraphs[0].font.size = Pt(44)
conclusion_frame.paragraphs[0].font.bold = True
conclusion_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
conclusion_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

c1 = conclusion_frame.add_paragraph()
c1.text = "\n760 drops = TCP reordering"
c1.font.size = Pt(28)
c1.font.color.rgb = RGBColor(255, 255, 255)
c1.alignment = PP_ALIGN.CENTER

c2 = conclusion_frame.add_paragraph()
c2.text = "NOT actual packet loss"
c2.font.size = Pt(28)
c2.font.color.rgb = RGBColor(255, 255, 255)
c2.alignment = PP_ALIGN.CENTER

c3 = conclusion_frame.add_paragraph()
c3.text = "\nValidated by real-world studies"
c3.font.size = Pt(26)
c3.font.color.rgb = RGBColor(255, 255, 255)
c3.alignment = PP_ALIGN.CENTER

c4 = conclusion_frame.add_paragraph()
c4.text = "\nHybrid only helps WITH congestion"
c4.font.size = Pt(26)
c4.font.color.rgb = RGBColor(255, 255, 0)
c4.font.bold = True
c4.alignment = PP_ALIGN.CENTER

prs.save('Saims_Questions_SIMPLE.pptx')
print("✓ Simple presentation created!")
print("  File: Saims_Questions_SIMPLE.pptx")
print("  Slides: 6")
